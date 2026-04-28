"""Generate the AI Network Guardian presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme colors ──────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x11, 0x1A)
BG_CARD  = RGBColor(0x1A, 0x1D, 0x2E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x8B, 0x8F, 0xA8)
ACCENT   = RGBColor(0x6C, 0x63, 0xFF)
GREEN    = RGBColor(0x00, 0xE6, 0x76)
RED      = RGBColor(0xFF, 0x4D, 0x6A)
YELLOW   = RGBColor(0xFF, 0xD9, 0x3D)
CYAN     = RGBColor(0x00, 0xD4, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text,
             size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
             font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_para(tf, text, size=18, color=WHITE, bold=False, space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.space_before = Pt(space_before)
    return p


def add_card(slide, left, top, width, height, color=BG_CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_accent_line(slide, top=1.6):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(top), Inches(2), Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def slide_title(slide, title, subtitle=None):
    set_bg(slide)
    add_text(slide, 0.8, 0.5, 11, 0.9, title, size=36, bold=True, color=WHITE)
    add_accent_line(slide)
    if subtitle:
        add_text(slide, 0.8, 1.8, 11, 0.6, subtitle, size=18, color=MUTED)


# ══════════════════════════════════════════════════════════════
#  SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(s)

add_text(s, 0, 2.0, 13.333, 1.2,
         "AI NETWORK GUARDIAN", size=48, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, 0, 3.2, 13.333, 0.7,
         "An AI-Assisted TCP/IP Network Diagnosis System", size=24,
         color=ACCENT, align=PP_ALIGN.CENTER)

# accent line centered
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(5.166), Inches(4.2), Inches(3), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text(s, 0, 4.6, 13.333, 0.5,
         "Hamide Sila AKDAN  &  Bekir Sadik ALTUNKAYA", size=20,
         color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, 0, 5.2, 13.333, 0.5,
         "Mobile Technologies — Spring 2026", size=16,
         color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(s, "The Problem")

problems = [
    ("Unknown Devices", "Who is connected to my network? Is someone unauthorized?"),
    ("Website Trust", "Is this website safe? Is it a phishing or scam page?"),
    ("Performance Issues", "Why is my internet slow? What is causing lag?"),
]
for i, (title, desc) in enumerate(problems):
    x = 0.8 + i * 4.0
    add_card(s, x, 2.5, 3.6, 2.2)
    add_text(s, x + 0.3, 2.7, 3.0, 0.5, title, size=22, bold=True, color=ACCENT)
    add_text(s, x + 0.3, 3.3, 3.0, 1.2, desc, size=16, color=MUTED)

add_text(s, 0.8, 5.2, 11, 0.8,
         "Existing tools give raw, technical outputs that non-expert users cannot interpret.",
         size=18, color=YELLOW)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3 — Our Solution
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(s, "Our Solution")

tf = add_text(s, 0.8, 2.3, 11, 4.0,
              "AI Network Guardian transforms raw network data into clear, "
              "actionable diagnostics that anyone can understand.", size=20, color=WHITE)
add_para(tf, "", size=10)
add_para(tf, "Every diagnosis includes:", size=20, color=ACCENT, bold=True)
add_para(tf, "   TCP/IP Layer mapping — which layer is responsible", size=18, color=WHITE)
add_para(tf, "   Confidence score — how certain the system is (0-100%)", size=18, color=WHITE)
add_para(tf, "   Supporting evidence — specific data points from measurements", size=18, color=WHITE)
add_para(tf, "   Plain-English explanation — no jargon, no black boxes", size=18, color=WHITE)
add_para(tf, "   Actionable recommendation — what to do about it", size=18, color=WHITE)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4 — Architecture Overview
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(s, "System Architecture")

# Data flow boxes
components = [
    (1.5, 2.3, "User / Browser",      MUTED),
    (1.5, 3.3, "Flask API Gateway",    ACCENT),
    (1.5, 4.3, "Analysis Modules",     CYAN),
    (1.5, 5.3, "AI Reasoning Core",    GREEN),
    (1.5, 6.3, "SQLite Time-Series DB",YELLOW),
]
for x, y, label, color in components:
    add_card(s, x, y, 3.5, 0.7)
    add_text(s, x + 0.2, y + 0.1, 3.1, 0.5, label, size=16, bold=True, color=color)

# Arrows (simple text)
for y in [3.0, 4.0, 5.0, 6.0]:
    add_text(s, 2.8, y, 0.5, 0.3, "v", size=20, color=MUTED, align=PP_ALIGN.CENTER)

# Right side — module details
add_card(s, 6.5, 2.3, 5.8, 4.7)
tf = add_text(s, 6.8, 2.5, 5.2, 0.5, "Three Analysis Modules", size=22, bold=True, color=ACCENT)

modules_info = [
    ("Network Detective", "L1/L2", "ARP, MAC, Ping Sweep"),
    ("Security Hunter",   "L7",    "HTTPS Certs, WHOIS, Phishing Detection"),
    ("Performance Monitor","L3/L4","TCP Latency, Connections, DNS"),
    ("Connection Test",   "L1-L4", "TCP Handshake Layer Analysis"),
]
for i, (name, layer, desc) in enumerate(modules_info):
    y = 3.2 + i * 0.9
    add_text(s, 6.8, y, 2.5, 0.4, name, size=17, bold=True, color=WHITE)
    add_text(s, 9.5, y, 1.0, 0.4, layer, size=15, bold=True, color=CYAN)
    add_text(s, 6.8, y + 0.35, 5.0, 0.4, desc, size=14, color=MUTED)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5 — TCP/IP Layer Mapping
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(s, "TCP/IP Layer Mapping")

layers = [
    ("Application (L7)",  "HTTP, HTTPS, DNS, WHOIS",       "Security Hunter",      RGBColor(0xE0, 0x60, 0xFF)),
    ("Transport (L4)",    "TCP, UDP — ports, sessions",     "Performance Monitor",  RGBColor(0x00, 0xD4, 0xFF)),
    ("Network (L3)",      "IP, ICMP — routing, latency",   "Performance Monitor",  RGBColor(0x00, 0xE6, 0x76)),
    ("Data Link (L2)",    "Ethernet, ARP, MAC",             "Network Detective",    RGBColor(0xFF, 0xD9, 0x3D)),
]
for i, (layer, protocols, module, color) in enumerate(layers):
    y = 2.3 + i * 1.15
    add_card(s, 0.8, y, 11.5, 0.95)
    add_text(s, 1.1, y + 0.15, 3.0, 0.5, layer, size=20, bold=True, color=color)
    add_text(s, 4.3, y + 0.15, 4.0, 0.5, protocols, size=16, color=MUTED)
    add_text(s, 8.8, y + 0.15, 3.2, 0.5, module, size=16, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6 — AI Reasoning Engine
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(s, "AI Reasoning Engine")

# Left — Primary
add_card(s, 0.8, 2.3, 5.5, 4.0)
add_text(s, 1.1, 2.5, 5.0, 0.5, "Primary: Google Gemini 2.5", size=22, bold=True, color=GREEN)
tf = add_text(s, 1.1, 3.2, 5.0, 3.0,
              "Plain HTTP REST API (no SDK dependency)", size=16, color=WHITE)
add_para(tf, "Flash-first, auto-fallback to Pro on rate limit", size=16, color=WHITE)
add_para(tf, "Structured system prompt with TCP/IP context", size=16, color=WHITE)
add_para(tf, "Institutional domain awareness (.edu, .gov)", size=16, color=WHITE)
add_para(tf, "Phishing infrastructure pattern detection", size=16, color=WHITE)
add_para(tf, "Human-readable Website Overview generation", size=16, color=WHITE)

# Right — Fallback
add_card(s, 6.8, 2.3, 5.5, 4.0)
add_text(s, 7.1, 2.5, 5.0, 0.5, "Fallback: Rule-Based Engine", size=22, bold=True, color=YELLOW)
tf = add_text(s, 7.1, 3.2, 5.0, 3.0,
              "Deterministic rules when API is unavailable", size=16, color=WHITE)
add_para(tf, "Same Diagnosis dataclass output format", size=16, color=WHITE)
add_para(tf, "Threshold-based severity classification", size=16, color=WHITE)
add_para(tf, "Zero external dependencies for offline use", size=16, color=WHITE)
add_para(tf, "Consistent UX regardless of AI availability", size=16, color=WHITE)

add_text(s, 0.8, 6.5, 11, 0.5,
         "Both engines produce identical output structure — the UI never knows which one ran.",
         size=16, color=MUTED)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7 — Module 1: Network Detective
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(s, "Module 1: Network Detective", "Layer 1/2 — ARP, MAC, ICMP")

tf = add_text(s, 0.8, 2.3, 5.5, 4.5, "", size=16, color=WHITE)
add_para(tf, "How it works:", size=20, bold=True, color=ACCENT, space_before=0)
add_para(tf, "1. Reads ARP table to discover local devices", size=17, color=WHITE)
add_para(tf, "2. Pings each device to verify responsiveness", size=17, color=WHITE)
add_para(tf, "3. Looks up MAC OUI prefix for vendor identification", size=17, color=WHITE)
add_para(tf, "4. Detects randomized/private MACs (iOS/Android)", size=17, color=WHITE)
add_para(tf, "5. Identifies gateway and network topology", size=17, color=WHITE)
add_para(tf, "", size=10)
add_para(tf, "AI analyzes the scan for:", size=20, bold=True, color=ACCENT)
add_para(tf, "  Unknown/unauthorized devices", size=17, color=WHITE)
add_para(tf, "  Non-responsive hosts (firewalled)", size=17, color=WHITE)
add_para(tf, "  Abnormal device counts", size=17, color=WHITE)

add_card(s, 7.0, 2.3, 5.5, 4.5)
add_text(s, 7.3, 2.5, 5.0, 0.5, "Key Protocols", size=20, bold=True, color=CYAN)
protocols = [
    ("ARP", "Address Resolution Protocol — maps IP to MAC at L2"),
    ("ICMP", "Ping sweep — verifies host reachability at L3"),
    ("MAC OUI", "First 3 bytes identify the hardware manufacturer"),
]
for i, (proto, desc) in enumerate(protocols):
    y = 3.2 + i * 1.0
    add_text(s, 7.3, y, 1.5, 0.4, proto, size=18, bold=True, color=GREEN)
    add_text(s, 7.3, y + 0.4, 4.8, 0.5, desc, size=14, color=MUTED)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8 — Module 2: Security Hunter
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(s, "Module 2: Security Hunter", "Layer 7 — HTTPS, SSL/TLS, WHOIS, DNS")

tf = add_text(s, 0.8, 2.3, 5.5, 4.5, "", size=16, color=WHITE)
add_para(tf, "Multi-signal phishing detection:", size=20, bold=True, color=ACCENT, space_before=0)
add_para(tf, "  SSL certificate validation & metadata", size=17, color=WHITE)
add_para(tf, "  WHOIS domain registration analysis", size=17, color=WHITE)
add_para(tf, "  TLS version analysis (1.2 vs 1.3)", size=17, color=WHITE)
add_para(tf, "  Wildcard & multi-domain cert detection", size=17, color=WHITE)
add_para(tf, "  Government impersonation detection", size=17, color=WHITE)
add_para(tf, "  Scam bait keyword analysis", size=17, color=WHITE)
add_para(tf, "  Institutional domain trust (.edu, .gov)", size=17, color=WHITE)

add_card(s, 7.0, 2.3, 5.5, 4.5)
add_text(s, 7.3, 2.5, 5.0, 0.5, "Scoring System", size=20, bold=True, color=CYAN)
tf2 = add_text(s, 7.3, 3.2, 5.0, 3.5, "", size=16, color=WHITE)
add_para(tf2, "Risk signals add points (0-100)", size=17, color=RED, space_before=0)
add_para(tf2, "Trust signals subtract points", size=17, color=GREEN)
add_para(tf2, "Trust capped at 50% of risk", size=17, color=YELLOW)
add_para(tf2, "(prevents false negatives)", size=14, color=MUTED)
add_para(tf2, "", size=10)
add_para(tf2, "0-30:  SAFE", size=17, color=GREEN)
add_para(tf2, "31-60: CAUTION", size=17, color=YELLOW)
add_para(tf2, "61-100: DANGEROUS", size=17, color=RED)


# ══════════════════════════════════════════════════════════════
#  SLIDE 9 — Security Hunter Results
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(s, "Security Hunter: Real Results")

# Test cases
cases = [
    ("google.com",              "2/100",  "SAFE",      GREEN, "Premium CA, 28yr domain, valid cert"),
    ("gtu.edu.tr",              "8/100",  "SAFE",      GREEN, "TUBITAK gov CA, institutional .edu.tr"),
    ("phishing.filterdns.net",  "70/100", "CRITICAL",  RED,   "High-threat keyword, wildcard + free CA, subdomain impersonation"),
    ("govfreephone.us",         "80/100", "CRITICAL",  RED,   "Gov impersonation, scam bait, multi-domain cert, phishing infra"),
]

for i, (domain, score, level, color, reason) in enumerate(cases):
    y = 2.3 + i * 1.2
    add_card(s, 0.8, y, 11.5, 1.0)
    add_text(s, 1.1, y + 0.1, 3.5, 0.4, domain, size=18, bold=True, color=WHITE)
    add_text(s, 4.8, y + 0.1, 1.2, 0.4, score, size=18, bold=True, color=color)
    add_text(s, 6.2, y + 0.1, 1.5, 0.4, level, size=16, bold=True, color=color)
    add_text(s, 1.1, y + 0.5, 10.5, 0.4, reason, size=14, color=MUTED)


# ══════════════════════════════════════════════════════════════
#  SLIDE 10 — Module 3: Performance Monitor
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(s, "Module 3: Performance & Lag Monitor", "Layer 3/4 — IP, TCP, UDP, ICMP")

tf = add_text(s, 0.8, 2.3, 5.5, 4.5, "", size=16, color=WHITE)
add_para(tf, "Measurements:", size=20, bold=True, color=ACCENT, space_before=0)
add_para(tf, "  TCP connection latency (min/avg/max)", size=17, color=WHITE)
add_para(tf, "  Packet loss percentage", size=17, color=WHITE)
add_para(tf, "  Jitter (latency variance)", size=17, color=WHITE)
add_para(tf, "  Active TCP/UDP connection states", size=17, color=WHITE)
add_para(tf, "  DNS resolution time", size=17, color=WHITE)
add_para(tf, "  Top remote hosts by connection count", size=17, color=WHITE)
add_para(tf, "", size=10)
add_para(tf, "No ICMP / no root required", size=18, bold=True, color=GREEN)
add_para(tf, "Uses TCP SYN handshake for latency", size=16, color=MUTED)

add_card(s, 7.0, 2.3, 5.5, 4.5)
add_text(s, 7.3, 2.5, 5.0, 0.5, "AI Diagnoses", size=20, bold=True, color=CYAN)
tf2 = add_text(s, 7.3, 3.2, 5.0, 3.5, "", size=16, color=WHITE)
add_para(tf2, "Elevated latency → cause analysis", size=17, color=WHITE, space_before=0)
add_para(tf2, "Packet loss → severity assessment", size=17, color=WHITE)
add_para(tf2, "High jitter → VoIP/gaming impact", size=17, color=WHITE)
add_para(tf2, "Connection count → resource usage", size=17, color=WHITE)
add_para(tf2, "DNS failure → resolution troubleshooting", size=17, color=WHITE)
add_para(tf2, "", size=10)
add_para(tf2, "Each mapped to responsible TCP/IP layer", size=16, color=MUTED)


# ══════════════════════════════════════════════════════════════
#  SLIDE 11 — Module 4: Connection Test
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(s, "Module 4: Connection Test", "Layer 1-4 — Full Stack TCP Analysis")

statuses = [
    ("CONNECTED",   "SYN-ACK received, port open",           "L1-L4 all functional",    GREEN),
    ("REFUSED",     "TCP RST received, port closed",          "L1-L3 OK, L4 closed",     YELLOW),
    ("TIMEOUT",     "No response, firewall filtering",        "L4 filtered",             RED),
    ("UNREACHABLE", "No route to host, L3 failure",           "L2/L3 failure",           RED),
]

for i, (status, meaning, layers, color) in enumerate(statuses):
    y = 2.3 + i * 1.15
    add_card(s, 0.8, y, 11.5, 0.95)
    add_text(s, 1.1, y + 0.1, 2.5, 0.4, status, size=20, bold=True, color=color)
    add_text(s, 3.8, y + 0.1, 4.5, 0.4, meaning, size=16, color=WHITE)
    add_text(s, 8.8, y + 0.1, 3.2, 0.4, layers, size=16, color=MUTED)

add_text(s, 0.8, 6.5, 11, 0.5,
         "Each status maps to a specific layer failure — the AI explains what went wrong and how to fix it.",
         size=16, color=MUTED)


# ══════════════════════════════════════════════════════════════
#  SLIDE 12 — Data Persistence
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(s, "Time-Series Data & History")

tf = add_text(s, 0.8, 2.3, 11, 4.5, "", size=16, color=WHITE)
add_para(tf, "SQLite with WAL mode for concurrent reads", size=18, color=WHITE, space_before=0)
add_para(tf, "", size=8)
add_para(tf, "Every scan result is persisted with:", size=20, bold=True, color=ACCENT)
add_para(tf, "   Module identifier (detective, security, performance, connection)", size=17, color=WHITE)
add_para(tf, "   Unix timestamp for time-series ordering", size=17, color=WHITE)
add_para(tf, "   Raw measurement data (JSON)", size=17, color=WHITE)
add_para(tf, "   AI-generated diagnoses (JSON)", size=17, color=WHITE)
add_para(tf, "", size=8)
add_para(tf, "API Endpoints:", size=20, bold=True, color=ACCENT)
add_para(tf, "   GET /api/history/<module> — recent scan results", size=17, color=CYAN)
add_para(tf, "   GET /api/trend/<module>?hours=24 — time-windowed data", size=17, color=CYAN)
add_para(tf, "", size=8)
add_para(tf, "Purpose: Distinguish temporary anomalies from persistent issues", size=18, color=YELLOW, bold=True)


# ══════════════════════════════════════════════════════════════
#  SLIDE 13 — Tech Stack
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(s, "Technology Stack")

stack = [
    ("Backend",   "Python 3 + Flask",              "REST API, routing, module orchestration"),
    ("AI",        "Google Gemini 2.5 (Flash/Pro)",  "Plain HTTP REST, no SDK, auto-fallback"),
    ("Database",  "SQLite + WAL",                   "Time-series storage, indexed queries"),
    ("Frontend",  "Vanilla HTML/CSS/JS",            "Dark theme, tab-based SPA, no framework"),
    ("Security",  "python-whois, ssl, socket",      "Certificate parsing, WHOIS lookup"),
    ("Network",   "arp, netstat, socket",            "ARP table, connections, TCP probes"),
]

for i, (category, tech, desc) in enumerate(stack):
    y = 2.3 + i * 0.8
    add_card(s, 0.8, y, 11.5, 0.65)
    add_text(s, 1.1, y + 0.1, 2.0, 0.4, category, size=16, bold=True, color=ACCENT)
    add_text(s, 3.3, y + 0.1, 3.5, 0.4, tech, size=16, bold=True, color=WHITE)
    add_text(s, 7.0, y + 0.1, 5.0, 0.4, desc, size=14, color=MUTED)


# ══════════════════════════════════════════════════════════════
#  SLIDE 14 — Live Demo Slide
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

add_text(s, 0, 2.5, 13.333, 1.0,
         "LIVE DEMO", size=54, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER)

shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(5.166), Inches(3.8), Inches(3), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text(s, 0, 4.2, 13.333, 0.5,
         "http://127.0.0.1:5001", size=24, color=CYAN,
         align=PP_ALIGN.CENTER)

add_text(s, 0, 5.2, 13.333, 0.8,
         "Network Scan  |  Security Analysis  |  Performance Diagnostics  |  Connection Test",
         size=18, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 15 — Conclusion
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_title(s, "Conclusion")

achievements = [
    "Three integrated modules covering all TCP/IP layers (L2-L7)",
    "AI-powered diagnostics with Gemini 2.5 + deterministic fallback",
    "Every diagnosis: layer-mapped, confidence-scored, evidence-backed",
    "Advanced phishing detection: gov impersonation, scam bait, multi-domain certs",
    "Institutional domain awareness: .edu, .gov, government CAs",
    "Time-series persistence for anomaly vs. persistent issue distinction",
    "No root/admin privileges required — runs on any machine",
    "Clean, accessible dark-theme UI for non-technical users",
]

for i, item in enumerate(achievements):
    y = 2.3 + i * 0.55
    add_text(s, 1.1, y, 0.4, 0.4, "✓", size=18, bold=True, color=GREEN)
    add_text(s, 1.6, y, 10.5, 0.4, item, size=17, color=WHITE)

add_text(s, 0, 6.5, 13.333, 0.5,
         "Thank you — Questions?", size=28, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────
out = "AI_Network_Guardian_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")